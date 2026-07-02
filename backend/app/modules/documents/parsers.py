"""Factory Method pattern: mỗi FileType có một creator riêng, factory method
`_create_parser` quyết định trả về parser cụ thể nào -> caller chỉ biết
interface `DocumentParser`, không biết chi tiết PDF/DOCX/PPTX."""
import io
from abc import ABC, abstractmethod

from app.modules.documents.models import FileType


class DocumentParser(ABC):
    """Product: interface chung cho mọi parser tài liệu."""

    @abstractmethod
    def parse(self, content: bytes) -> list[tuple[str, int]]:
        """Trả về danh sách (text, page_number)."""


class PdfParser(DocumentParser):
    def parse(self, content: bytes) -> list[tuple[str, int]]:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        pages = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append((text, i))
        return pages


class DocxParser(DocumentParser):
    def parse(self, content: bytes) -> list[tuple[str, int]]:
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(content))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        return [(text, 1)] if text.strip() else []


class PptxParser(DocumentParser):
    def parse(self, content: bytes) -> list[tuple[str, int]]:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        pages = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = [
                shape.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ]
            if parts:
                pages.append(("\n".join(parts), i))
        return pages


# Factory method: ánh xạ FileType -> creator tạo đúng parser cụ thể.
_PARSER_FACTORIES: dict[FileType, type[DocumentParser]] = {
    FileType.PDF: PdfParser,
    FileType.DOCX: DocxParser,
    FileType.PPTX: PptxParser,
}


def create_parser(file_type: FileType) -> DocumentParser:
    """Factory method: nhận FileType, trả về instance DocumentParser phù hợp."""
    return _PARSER_FACTORIES[file_type]()


def detect_file_type(filename: str) -> FileType:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return FileType.PDF
    if lower.endswith(".docx"):
        return FileType.DOCX
    if lower.endswith(".pptx"):
        return FileType.PPTX
    raise ValueError("Định dạng không hỗ trợ (chỉ PDF, DOCX, PPTX)")


# FR-LEC-01 / §11.3 / §12: chỉ chấp nhận PDF, DOCX, PPTX — validate cả MIME
# type lẫn extension ở server side.
_ALLOWED_MIME: dict[FileType, set[str]] = {
    FileType.PDF: {"application/pdf"},
    FileType.DOCX: {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    },
    FileType.PPTX: {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    },
}

# MIME chung chung mà nhiều client gửi khi không nhận diện được — khi đó tin
# vào extension đã được kiểm ở detect_file_type.
_GENERIC_MIME = {"", "application/octet-stream", "binary/octet-stream"}


def validate_mime(file_type: FileType, content_type: str | None) -> None:
    if (content_type or "") in _GENERIC_MIME:
        return
    if content_type not in _ALLOWED_MIME[file_type]:
        raise ValueError(
            f"Nội dung tệp ({content_type}) không khớp định dạng "
            f"{file_type.value}"
        )


def parse(content: bytes, file_type: FileType) -> list[tuple[str, int]]:
    return create_parser(file_type).parse(content)
